import sys
import os

try:
    import pptx
except ImportError:
    import subprocess
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_pitch_deck(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette - Dark Theme matching Insights Copilot UI
    BG_COLOR = RGBColor(15, 23, 42)        # #0F172A Dark Slate
    CARD_BG = RGBColor(30, 41, 59)        # #1E293B Card Background
    ACCENT_COLOR = RGBColor(99, 102, 241)  # #6366F1 Indigo Accent
    CYAN_ACCENT = RGBColor(56, 189, 248)   # #38BDF8 Cyan Accent
    TEXT_WHITE = RGBColor(248, 250, 252)   # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8 Muted Grey
    BORDER_COLOR = RGBColor(51, 65, 85)   # #334155

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def create_card(slide, left, top, width, height, bg_rgb=CARD_BG, border_rgb=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_rgb
        if border_rgb:
            shape.line.color.rgb = border_rgb
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def add_header(slide, title_text, category_text=None):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        if category_text:
            p0 = tf.paragraphs[0]
            p0.text = category_text.upper()
            p0.font.size = Pt(11)
            p0.font.bold = True
            p0.font.color.rgb = CYAN_ACCENT
            p0.font.name = "Segoe UI"
            
            p1 = tf.add_paragraph()
        else:
            p1 = tf.paragraphs[0]
            
        p1.text = title_text
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.font.name = "Segoe UI"

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    create_card(slide1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5), bg_rgb=CARD_BG, border_rgb=ACCENT_COLOR)

    tbox1 = slide1.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.333), Inches(3.5))
    tf1 = tbox1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "INSIGHTS COPILOT"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "Autonomous AI Co-Founder"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_WHITE
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "\"An AI co-founder that scores your startup idea and turns it into a milestone plan you'll actually follow.\""
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Segoe UI"
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)

    # ==========================================
    # SLIDE 2: Problem
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Early-Stage Founder Bottleneck", "THE PROBLEM")

    problems = [
        ("Lack of Honest Evaluation", [
            "Builders create endless ideas",
            "No objective validation metric",
            "Biased feedback from peers",
            "High chance of building dead-ends"
        ]),
        ("No Execution Structure", [
            "Lack concrete milestone roadmaps",
            "Overwhelmed by initial scope",
            "Difficulty maintaining momentum",
            "High rate of abandoned projects"
        ]),
        ("Flaws of Generic AI Chat", [
            "Chatbots brainstorm without scoring",
            "Zero progress or milestone tracking",
            "No proactive follow-ups or nudges",
            "Context lost after session ends"
        ])
    ]

    col_width = Inches(3.64)
    for i, (title, bullets) in enumerate(problems):
        left = Inches(0.8 + i * 3.95)
        create_card(slide2, left, Inches(1.8), col_width, Inches(5.0))
        
        tb = slide2.shapes.add_textbox(left + Inches(0.3), Inches(2.1), col_width - Inches(0.6), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        p.font.name = "Segoe UI"
        
        for b in bullets:
            bp = tf.add_paragraph()
            bp.text = "• " + b
            bp.font.size = Pt(14)
            bp.font.color.rgb = TEXT_WHITE
            bp.font.name = "Segoe UI"
            bp.space_before = Pt(12)

    # ==========================================
    # SLIDE 3: Solution
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Autonomous End-to-End Execution Pipeline", "THE SOLUTION")

    pipeline_steps = [
        ("1. Input", "Raw Startup Idea", "Founder submits concept"),
        ("2. Research & Score", "Planner/Critic Pipeline", "Evaluates feasibility & impact"),
        ("3. Roadmap", "Milestone Generation", "Actionable step-by-step tasks"),
        ("4. Dashboard", "Builder Workspace", "Tracks real-time score & state"),
        ("5. Follow-up", "Telegram Nudge Bot", "Automated proactive reminders")
    ]

    step_w = Inches(2.18)
    for i, (num, name, desc) in enumerate(pipeline_steps):
        left = Inches(0.8 + i * 2.36)
        create_card(slide3, left, Inches(2.2), step_w, Inches(4.4), border_rgb=ACCENT_COLOR if i%2==1 else BORDER_COLOR)
        
        tb = slide3.shapes.add_textbox(left + Inches(0.2), Inches(2.4), step_w - Inches(0.4), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = num
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = CYAN_ACCENT
        p0.font.name = "Segoe UI"
        
        p1 = tf.add_paragraph()
        p1.text = name
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.font.name = "Segoe UI"
        p1.space_before = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_MUTED
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(12)

    # ==========================================
    # SLIDE 4: Product Demo (Builder Dashboard)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Builder Dashboard & Founder Profile", "PRODUCT DEMO")

    create_card(slide4, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.0))

    metrics = [
        ("ACTIVE WORKSPACES", "4 Ideas"),
        ("COMPOSITE SCORE", "82 / 100"),
        ("MILESTONES DONE", "12 Completed"),
        ("TELEGRAM STATUS", "Connected")
    ]
    for i, (m_label, m_val) in enumerate(metrics):
        m_left = Inches(1.1 + i * 2.85)
        create_card(slide4, m_left, Inches(2.0), Inches(2.6), Inches(1.0), bg_rgb=BG_COLOR)
        tb = slide4.shapes.add_textbox(m_left + Inches(0.1), Inches(2.1), Inches(2.4), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_label
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.font.name = "Segoe UI"
        p2 = tf.add_paragraph()
        p2.text = m_val
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = CYAN_ACCENT
        p2.font.name = "Segoe UI"

    cards = [
        ("AI Code Reviewer", "Score: 88 | Feasibility: High", "Milestone 3/5: Integration Tests"),
        ("Market Analytics SaaS", "Score: 74 | Feasibility: Med", "Milestone 1/4: Customer Survey"),
        ("DevOps Automator", "Score: 81 | Feasibility: High", "Milestone 4/4: Deploy Staging")
    ]
    for i, (w_name, w_score, w_ms) in enumerate(cards):
        c_left = Inches(1.1 + i * 3.8)
        create_card(slide4, c_left, Inches(3.3), Inches(3.5), Inches(3.1), bg_rgb=BG_COLOR)
        tb = slide4.shapes.add_textbox(c_left + Inches(0.2), Inches(3.5), Inches(3.1), Inches(2.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = w_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = "Segoe UI"
        
        p2 = tf.add_paragraph()
        p2.text = w_score
        p2.font.size = Pt(13)
        p2.font.color.rgb = CYAN_ACCENT
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(8)
        
        p3 = tf.add_paragraph()
        p3.text = w_ms
        p3.font.size = Pt(12)
        p3.font.color.rgb = TEXT_MUTED
        p3.font.name = "Segoe UI"
        p3.space_before = Pt(8)

    # ==========================================
    # SLIDE 5: Product Demo (Needs Attention & Nudges)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Progress Tracking & Proactive Nudges", "PRODUCT DEMO")

    create_card(slide5, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.0))
    tb_l = slide5.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Needs Attention Dashboard View"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.font.name = "Segoe UI"

    bullets_l = [
        "Identifies stale workspace milestones",
        "Flags stagnant ideas needing action",
        "Tracks time elapsed since last check-in",
        "Direct link to resolve milestone blockers"
    ]
    for b in bullets_l:
        bp = tf_l.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.font.name = "Segoe UI"
        bp.space_before = Pt(14)

    create_card(slide5, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), border_rgb=ACCENT_COLOR)
    tb_r = slide5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "Telegram Nudge Loop Integration"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.font.name = "Segoe UI"

    bullets_r = [
        "Proactive milestone reminder alerts",
        "Weekly automated execution digest",
        "Direct /done commands via chat",
        "Real-time sync back to Web App"
    ]
    for b in bullets_r:
        bp = tf_r.add_paragraph()
        bp.text = "• " + b
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.font.name = "Segoe UI"
        bp.space_before = Pt(14)

    # ==========================================
    # SLIDE 6: Key Differentiation
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "System of Record & Multi-Provider Architecture", "DIFFERENTIATION")

    diff_blocks = [
        ("System of Record over Time", [
            "One-shot tools score once & vanish",
            "Insights Copilot tracks full lifecycle",
            "Maintains long-term idea history",
            "Enables persistent Telegram nudge loop"
        ]),
        ("Multi-Provider Reliability Architecture", [
            "Groq (Llama 3.3 70B) primary engine",
            "Gemini 2.5 Flash automatic fallback",
            "Zero downtime during rate limits",
            "Engineering resilience over UI tricks"
        ])
    ]

    for i, (title, bullets) in enumerate(diff_blocks):
        left = Inches(0.8 + i * 5.95)
        create_card(slide6, left, Inches(1.8), Inches(5.7), Inches(5.0))
        
        tb = slide6.shapes.add_textbox(left + Inches(0.4), Inches(2.1), Inches(4.9), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        p.font.name = "Segoe UI"
        
        for b in bullets:
            bp = tf.add_paragraph()
            bp.text = "• " + b
            bp.font.size = Pt(14)
            bp.font.color.rgb = TEXT_WHITE
            bp.font.name = "Segoe UI"
            bp.space_before = Pt(16)

    # ==========================================
    # SLIDE 7: Target Users
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Tailored for Builders & Institutions", "TARGET USERS")

    targets = [
        ("Primary: Early-Stage Builders", [
            "Hackathon participants & students",
            "Side-project developers",
            "Need objective score validation",
            "Require structured milestone execution"
        ]),
        ("Secondary: Innovation Institutions", [
            "Startup incubators & accelerators",
            "College entrepreneurship cells",
            "Batch idea triage & evaluation",
            "Cohort-wide founder progress tracking"
        ])
    ]

    for i, (title, bullets) in enumerate(targets):
        left = Inches(0.8 + i * 5.95)
        create_card(slide7, left, Inches(1.8), Inches(5.7), Inches(5.0))
        
        tb = slide7.shapes.add_textbox(left + Inches(0.4), Inches(2.1), Inches(4.9), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        p.font.name = "Segoe UI"
        
        for b in bullets:
            bp = tf.add_paragraph()
            bp.text = "• " + b
            bp.font.size = Pt(14)
            bp.font.color.rgb = TEXT_WHITE
            bp.font.name = "Segoe UI"
            bp.space_before = Pt(16)

    # ==========================================
    # SLIDE 8: Business Model
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Freemium SaaS & Institution Licensing", "BUSINESS MODEL")

    table_shape = slide8.shapes.add_table(4, 4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.6))
    table = table_shape.table

    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(4.8)
    table.columns[3].width = Inches(2.8)

    headers = ["TIER", "PRICE", "CORE FEATURES", "TARGET SEGMENT"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = "Segoe UI"

    rows_data = [
        ("Free Tier", "₹0", "3 ideas/mo, basic scoring, capped history, no Telegram", "Hobbyists & Curious"),
        ("Builder Tier", "Low Monthly", "Unlimited ideas, full milestones, Telegram alerts, full history", "Active Solo Builders"),
        ("Cohort / Inst.", "Per-seat / Flat", "Everything in Builder + Mentor dashboard & cohort reports", "Incubators & E-Cells")
    ]

    for i, r_data in enumerate(rows_data):
        for j, val in enumerate(r_data):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_WHITE
            p.font.name = "Segoe UI"

    create_card(slide8, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2))
    tb_sec = slide8.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.0))
    tf_sec = tb_sec.text_frame
    tf_sec.word_wrap = True
    p = tf_sec.paragraphs[0]
    p.text = "Secondary Revenue Lever: Usage-Based BYO-Key / High-Volume Add-on"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.font.name = "Segoe UI"
    p2 = tf_sec.add_paragraph()
    p2.text = "• Enabled by existing multi-provider LLM routing architecture (Groq / Gemini / OpenAI keys)."
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 9: Roadmap & Honest Gaps
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Honest Gaps & Future Roadmap", "TRANSPARENCY")

    gaps = [
        ("Current Status (Built & Working)", [
            "Four-agent Planner/Critic pipeline",
            "FastAPI backend & PostgreSQL/Supabase DB",
            "React Builder Dashboard UI",
            "Telegram bot scheduler & handler"
        ]),
        ("Next Development Phase (Roadmap)", [
            "Multi-user Authentication (JWT/OAuth)",
            "Stripe/Razorpay billing integration",
            "Mentor aggregate cohort analytics view",
            "Custom webhooks & Discord integration"
        ])
    ]

    for i, (title, bullets) in enumerate(gaps):
        left = Inches(0.8 + i * 5.95)
        create_card(slide9, left, Inches(1.8), Inches(5.7), Inches(5.0))
        
        tb = slide9.shapes.add_textbox(left + Inches(0.4), Inches(2.1), Inches(4.9), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        p.font.name = "Segoe UI"
        
        for b in bullets:
            bp = tf.add_paragraph()
            bp.text = "• " + b
            bp.font.size = Pt(14)
            bp.font.color.rgb = TEXT_WHITE
            bp.font.name = "Segoe UI"
            bp.space_before = Pt(16)

    # ==========================================
    # SLIDE 10: Ask & Close Slide
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)

    create_card(slide10, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1), bg_rgb=CARD_BG, border_rgb=ACCENT_COLOR)

    tbox10 = slide10.shapes.add_textbox(Inches(2.0), Inches(1.5), Inches(9.333), Inches(4.5))
    tf10 = tbox10.text_frame
    tf10.word_wrap = True

    p = tf10.paragraphs[0]
    p.text = "JOIN US IN BUILDING THE FUTURE OF CO-FOUNDERS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    p_ask = tf10.add_paragraph()
    p_ask.text = "The Ask & Next Steps"
    p_ask.font.size = Pt(18)
    p_ask.font.bold = True
    p_ask.font.color.rgb = TEXT_WHITE
    p_ask.font.name = "Segoe UI"
    p_ask.space_before = Pt(20)

    ask_items = [
        "• Seeking Pre-Seed / Incubator Grant funding for API infra",
        "• Looking for Pilot Partners (Incubators, Hackathons, E-Cells)",
        "• Open for Technical Mentorship & Growth Advisors"
    ]
    for item in ask_items:
        bp = tf10.add_paragraph()
        bp.text = item
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.font.name = "Segoe UI"
        bp.space_before = Pt(8)

    p_contact = tf10.add_paragraph()
    p_contact.text = "Contact / Repository: github.com/Trusha25/insights-copilot"
    p_contact.font.size = Pt(14)
    p_contact.font.bold = True
    p_contact.font.color.rgb = CYAN_ACCENT
    p_contact.font.name = "Segoe UI"
    p_contact.alignment = PP_ALIGN.CENTER
    p_contact.space_before = Pt(24)

    prs.save(output_path)
    print(f"Pitch Deck successfully generated at: {output_path}")

if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_file = os.path.join(out_dir, "Insights_Copilot_Pitch_Deck.pptx")
    build_pitch_deck(out_file)
